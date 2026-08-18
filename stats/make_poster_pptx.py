"""
make_poster_pptx.py

Build an A0-portrait (84.1 x 118.9 cm) PowerPoint poster TEMPLATE that lays the
three poster figures out so they fit well, with editable placeholder text boxes
for the usual sections. The two wide landscape figures (Fig 2 fronto-striatal /
language FC, Fig 3 network size) stack full-width; the compact Fig 1 box plot
tucks in beside Fig 3.

Figures are embedded as high-resolution PNGs (rendered from the poster SVGs).
Re-run stats/make_poster_figures.py + the rsvg-convert PNG step first if the
figures changed. Everything (title, body text, captions) is editable in
PowerPoint; colours/fonts match the figures (Arial; salience-blue / language-red
accents).

Output: poster/poster_A0_template.pptx

Run:  .venv/bin/python stats/make_poster_pptx.py

## Author: Han Wang
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT_DIR = "/home/hanwang/Apps/Programming/matlab-proj/PFM_MSHBM_MHVerbGen"
POSTER = f"{PROJECT_DIR}/poster"
OUT = f"{POSTER}/poster_A0_template.pptx"

# ---- palette (matches the figures) -------------------------------------
BLUE = RGBColor(0x1F, 0x4E, 0x79)     # salience / primary theme
RED = RGBColor(0xA0, 0x20, 0x20)      # language accent
INK = RGBColor(0x22, 0x22, 0x22)      # body text
MUTE = RGBColor(0x55, 0x55, 0x55)     # captions
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xEC, 0xF1, 0xF7)     # pale header fill for text sections
FONT = "Arial"

# ---- A0 portrait canvas + layout grid (cm) -----------------------------
SW, SH = 84.1, 118.9
MX = 3.0                              # side margin
CW = SW - 2 * MX                      # content width = 78.1

prs = Presentation()
prs.slide_width = Cm(SW)
prs.slide_height = Cm(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
shapes = slide.shapes


def _set(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def rect(x, y, w, h, fill=None, line=None):
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp


def textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.2); tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.1); tf.margin_bottom = Cm(0.1)
    return tf


def para(tf, text, size, bold=False, color=INK, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, bullet=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = ("•  " + text) if bullet else text
    _set(r, size, bold=bold, color=color, italic=italic)
    return p


def section(x, y, w, h, title, accent=BLUE):
    """Header strip + empty body box; returns the body text_frame."""
    rect(x, y, w, 1.7, fill=accent)
    ht = textbox(x + 0.3, y, w - 0.6, 1.7, anchor=MSO_ANCHOR.MIDDLE)
    para(ht, title, 26, bold=True, color=WHITE, first=True, space_after=0)
    rect(x, y + 1.7, w, h - 1.7, fill=BAND, line=RGBColor(0xD5, 0xDE, 0xE8))
    body = textbox(x + 0.5, y + 2.0, w - 1.0, h - 2.3)
    return body


def figure(path, x, y, w):
    im = Image.open(path); ar = im.height / im.width
    h = w * ar
    shapes.add_picture(path, Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    return h


def caption(x, y, w, lead, rest, accent=BLUE):
    tf = textbox(x, y, w, 2.2)
    p = tf.paragraphs[0]; p.space_after = Pt(0)
    r = p.add_run(); r.text = lead + "  "; _set(r, 15, bold=True, color=accent)
    r = p.add_run(); r.text = rest; _set(r, 15, color=MUTE)


# ============================================================
# 1. Title band (full-bleed)
# ============================================================
rect(0, 0, SW, 11.0, fill=BLUE)
tf = textbox(MX, 0.7, CW, 7.2, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "[Poster title] Person-specific salience and language-network "
         "connectivity and emotional symptoms in Developmental Language Disorder",
     44, bold=True, color=WHITE, first=True, space_after=8)
tf2 = textbox(MX, 8.0, CW, 2.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf2, "[Author names]", 24, color=WHITE, first=True, space_after=2)
para(tf2, "[Affiliations · contact email]", 18, color=RGBColor(0xD6, 0xE2, 0xF0),
     italic=True)

# ============================================================
# 2. Intro row — Background & Aims | Methods
# ============================================================
bg = section(MX, 12.0, 37.0, 14.6, "Background & Aims")
para(bg, "Developmental Language Disorder (DLD) carries a markedly elevated risk "
         "of emotional difficulties; the brain basis is unclear.", 18, first=True)
para(bg, "The salience network and a left fronto-striatal / language circuit are "
         "candidate substrates for mood regulation.", 18)
para(bg, "Aim: test whether person-specific network size and connectivity relate "
         "to SDQ emotional symptoms, and whether that coupling differs by group.",
     18, bold=True)
para(bg, "[Edit: add 1–2 references or a hypothesis statement here.]", 15,
     italic=True, color=MUTE)

me = section(MX + 39.0, 12.0, CW - 39.0, 14.6, "Methods")
para(me, "Sample: n = 144 (DLD = 53, HSL = 27, TD = 64); icafix MS-HBM set n = 142.",
     18, first=True)
para(me, "Individualised cortical networks via MS-HBM; salience & language network "
         "size as % cortical surface.", 18)
para(me, "Fronto-striatal FC (9 tiles) and a left pars-opercularis–putamen language "
         "edge (CAB-NP), Fisher-z.", 18)
para(me, "Negative-binomial models of SDQ-emotional on group × predictor "
         "(TD reference); joint-interaction LR + Freedman–Lane permutation.", 18)
para(me, "[Edit: scanner/sequence, preprocessing, exclusions.]", 15,
     italic=True, color=MUTE)

# ============================================================
# 3. Results heading + figures
# ============================================================
rect(MX, 27.6, CW, 1.9, fill=BLUE)
rh = textbox(MX + 0.3, 27.6, CW - 0.6, 1.9, anchor=MSO_ANCHOR.MIDDLE)
para(rh, "Results", 26, bold=True, color=WHITE, first=True, space_after=0)

# Fig 2 (hero) — full-width, centred
f2w = 70.0
f2h = figure(f"{POSTER}/fig2_fc_scatter.png", (SW - f2w) / 2, 30.4, f2w)
caption((SW - f2w) / 2, 30.4 + f2h + 0.2, f2w, "Figure 2.",
        "FC → SDQ-emotional by group. Salience fronto-striatal grid (blue title) "
        "and the language frontal–putamen edge (red title); NB fits ± 95% CI.")

# Fig 3 (left) + Fig 1 (right), tops aligned
row_y = 30.4 + f2h + 3.0
f3w = 50.0
f3h = figure(f"{POSTER}/fig3_network_size.png", MX, row_y, f3w)
caption(MX, row_y + f3h + 0.2, f3w, "Figure 3.",
        "SDQ-emotional vs salience / language network size (NB mean ± 95% CI).")
f1w = 24.0
f1h = figure(f"{POSTER}/fig1_emotional_by_group.png", MX + f3w + 2.0, row_y, f1w)
caption(MX + f3w + 2.0, row_y + f1h + 0.2, f1w, "Figure 1.",
        "SDQ-emotional by group.")

# ============================================================
# 4. Conclusions + References
# ============================================================
concl_y = row_y + max(f3h, f1h) + 3.2
cc = section(MX, concl_y, CW, 11.0, "Conclusions")
para(cc, "[Key finding 1] Clinical groups show elevated emotional symptoms "
         "(large group main effects).", 19, first=True)
para(cc, "[Key finding 2] The FC → mood slope differs between DLD and TD for "
         "several fronto-striatal tiles and the language edge (interaction), "
         "though single-group slopes are mostly non-significant.", 19)
para(cc, "[Key finding 3 / take-home message — edit.]", 19, bold=True)

ref_y = concl_y + 11.0 + 0.6
rf = textbox(MX, ref_y, CW, 5.0)
para(rf, "References & Acknowledgements", 16, bold=True, color=BLUE, first=True,
     space_after=3)
para(rf, "[1] …   [2] …   [3] Gordon et al. 2021, Cereb Cortex.    "
         "Funding / acknowledgements: […].    Contact: [email].", 13, color=MUTE)

os.makedirs(POSTER, exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"  slide: {SW} x {SH} cm (A0 portrait)")
print(f"  Fig2 {f2w:.0f}x{f2h:.1f} cm | Fig3 {f3w:.0f}x{f3h:.1f} cm | "
      f"Fig1 {f1w:.0f}x{f1h:.1f} cm")
print(f"  content ends ~{ref_y + 5.0:.1f} cm (of {SH})")
